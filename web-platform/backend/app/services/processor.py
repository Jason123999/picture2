"""Core image processing services reused from the legacy CLI application."""

from __future__ import annotations

import io
import json
import logging
from dataclasses import dataclass
from pathlib import Path
from typing import Callable, Iterable, Sequence

import numpy as np
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt

_LOGGER = logging.getLogger(__name__)


@dataclass(slots=True)
class CropConfig:
    left: int
    top: int
    width: int
    height: int


@dataclass(slots=True)
class ImageProcessingSettings:
    dpi: int
    size_cm: float
    corner_radius_ratio: float
    label_config: dict[str, str]


class ImageProcessingError(Exception):
    """Raised when the processing pipeline fails."""


def cm_to_pixels(cm: float, dpi: int) -> int:
    return int(cm * dpi / 2.54)


def create_rounded_rectangle_mask(
    width: int, height: int, corner_radius_ratio: float
) -> Image.Image:
    mask = Image.new("L", (width, height), 0)
    draw = ImageDraw.Draw(mask)
    min_size = min(width, height)
    corner_radius = int(min_size * corner_radius_ratio)

    try:
        draw.rounded_rectangle(
            [0, 0, width - 1, height - 1], radius=corner_radius, fill=255
        )
    except AttributeError:
        draw.rectangle([corner_radius, 0, width - corner_radius, height], fill=255)
        draw.rectangle([0, corner_radius, width, height - corner_radius], fill=255)
        draw.ellipse([0, 0, corner_radius * 2, corner_radius * 2], fill=255)
        draw.ellipse([width - corner_radius * 2, 0, width, corner_radius * 2], fill=255)
        draw.ellipse([
            0,
            height - corner_radius * 2,
            corner_radius * 2,
            height,
        ], fill=255)
        draw.ellipse([
            width - corner_radius * 2,
            height - corner_radius * 2,
            width,
            height,
        ], fill=255)
    return mask


def crop_image_to_rounded_rectangle(
    image_path: Path | str,
    crop_coords: CropConfig,
    output_size_px: int,
    corner_radius_ratio: float,
) -> Image.Image:
    img = Image.open(image_path)
    left, top, width, height = (
        crop_coords.left,
        crop_coords.top,
        crop_coords.width,
        crop_coords.height,
    )
    cropped = img.crop((left, top, left + width, top + height))
    crop_width, crop_height = cropped.size

    if crop_width >= crop_height:
        output_width = output_size_px
        output_height = int(crop_height * output_size_px / crop_width)
    else:
        output_height = output_size_px
        output_width = int(crop_width * output_size_px / crop_height)

    resized = cropped.resize((output_width, output_height), Image.Resampling.LANCZOS)

    if corner_radius_ratio <= 0:
        if resized.mode != "RGBA":
            resized = resized.convert("RGBA")
        return resized

    mask = create_rounded_rectangle_mask(
        output_width, output_height, corner_radius_ratio
    )
    if resized.mode != "RGBA":
        resized = resized.convert("RGBA")

    resized_arr = np.array(resized, dtype=np.uint8)
    mask_arr = np.array(mask, dtype=np.uint8)
    if mask_arr.shape != (output_height, output_width):
        mask = mask.resize((output_width, output_height), Image.Resampling.NEAREST)
        mask_arr = np.array(mask, dtype=np.uint8)

    output_arr = np.zeros((output_height, output_width, 4), dtype=np.uint8)
    mask_bool = mask_arr > 128
    output_arr[mask_bool, :] = resized_arr[mask_bool, :]
    mask_outside = mask_arr == 0
    output_arr[mask_outside, :] = 0
    return Image.fromarray(output_arr, mode="RGBA")


def process_images(
    image_files: Sequence[Path | str],
    crop_config_provider: Callable[[Path | str], CropConfig],
    settings: ImageProcessingSettings,
) -> list[Image.Image]:
    processed: list[Image.Image] = []
    output_size_px = cm_to_pixels(settings.size_cm, settings.dpi)

    for index, image_file in enumerate(image_files, start=1):
        crop = crop_config_provider(image_file)
        if crop.width <= 0 or crop.height <= 0:
            _LOGGER.warning("Skip %s: invalid crop", image_file)
            continue
        image = crop_image_to_rounded_rectangle(
            image_file,
            crop,
            output_size_px,
            settings.corner_radius_ratio,
        )
        processed.append(image)
        _LOGGER.info("Processed %s/%s", index, len(image_files))

    return processed


def get_labels_for_image(filename: str, labels_config: dict[str, str]) -> dict[str, str]:
    return {
        "top_left": labels_config.get("top_left", ""),
        "top_right": labels_config.get("top_right", ""),
        "bottom_left": labels_config.get("bottom_left", ""),
        "bottom_right": labels_config.get("bottom_right", ""),
    }


def create_ppt(
    processed_images: Sequence[Image.Image],
    original_filenames: Sequence[str],
    settings: ImageProcessingSettings,
    columns: int,
    rows: int,
    column_spacing_cm: float,
    row_spacing_cm: float,
) -> Presentation:
    prs = Presentation()
    prs.slide_width = Cm(33.87)
    prs.slide_height = Cm(19.05)

    images_per_slide = columns * rows
    img_size_cm = settings.size_cm
    font_size = settings.label_config.get("font_size", 12)
    font_size_pt = Pt(font_size)

    label_color = settings.label_config.get("font_color", "#000000")
    font_color_rgb = RGBColor(
        int(label_color[1:3], 16),
        int(label_color[3:5], 16),
        int(label_color[5:7], 16),
    )

    total_width = columns * img_size_cm + (columns - 1) * column_spacing_cm
    total_height = rows * img_size_cm + (rows - 1) * row_spacing_cm
    slide_width_cm = 33.87
    slide_height_cm = 19.05
    start_x_cm = (slide_width_cm - total_width) / 2
    start_y_cm = (slide_height_cm - total_height) / 2

    slide = None
    for index, (img, filename) in enumerate(zip(processed_images, original_filenames)):
        if index % images_per_slide == 0:
            slide = prs.slides.add_slide(prs.slide_layouts[6])

        position = index % images_per_slide
        row = position // columns
        col = position % columns
        left_cm = start_x_cm + col * (img_size_cm + column_spacing_cm)
        top_cm = start_y_cm + row * (img_size_cm + row_spacing_cm)

        image_stream = io.BytesIO()
        img.save(image_stream, format="PNG")
        image_stream.seek(0)

        slide.shapes.add_picture(
            image_stream,
            Cm(left_cm),
            Cm(top_cm),
            width=Cm(img_size_cm),
            height=Cm(img_size_cm),
        )

        labels = get_labels_for_image(filename, settings.label_config)
        textbox_height_cm = 0.6
        textbox_width_cm = 1.2
        offset_cm = 0.1

        for position_key, label_key in [
            ("top_left", "top_left"),
            ("top_right", "top_right"),
            ("bottom_left", "bottom_left"),
            ("bottom_right", "bottom_right"),
        ]:
            label_text = labels.get(label_key)
            if not label_text:
                continue

            if "top" in position_key:
                y_pos = top_cm + offset_cm
            else:
                y_pos = top_cm + img_size_cm - textbox_height_cm - offset_cm

            if "left" in position_key:
                x_pos = left_cm + offset_cm
                align = PP_ALIGN.LEFT
            else:
                x_pos = left_cm + img_size_cm - textbox_width_cm - offset_cm
                align = PP_ALIGN.RIGHT

            textbox = slide.shapes.add_textbox(
                Cm(x_pos),
                Cm(y_pos),
                Cm(textbox_width_cm),
                Cm(textbox_height_cm),
            )
            text_frame = textbox.text_frame
            text_frame.text = label_text
            paragraph = text_frame.paragraphs[0]
            paragraph.font.size = font_size_pt
            paragraph.font.bold = True
            paragraph.font.color.rgb = font_color_rgb
            paragraph.alignment = align

            textbox.fill.solid()
            textbox.fill.fore_color.rgb = RGBColor(255, 255, 255)
            textbox.fill.transparency = 1.0
            textbox.line.fill.background()

    return prs


def load_config(config_path: Path) -> dict:
    with open(config_path, "r", encoding="utf-8") as file:
        return json.load(file)


def list_images(directory: Path, extensions: Iterable[str] | None = None) -> list[Path]:
    extensions = extensions or [".jpg", ".jpeg", ".tif", ".tiff", ".png"]
    files: list[Path] = []
    for ext in extensions:
        files.extend(sorted(directory.glob(f"*{ext}")))
    return files
